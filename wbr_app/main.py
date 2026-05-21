from fastapi import FastAPI, Query
from fastapi.responses import HTMLResponse, StreamingResponse
from google.cloud import bigquery
from datetime import datetime, timedelta
from io import BytesIO

app = FastAPI(title="WBR Report App")
client = bigquery.Client()
JOB_CONFIG = bigquery.QueryJobConfig(maximum_bytes_billed=50_000_000_000)

BQ_TABLE   = "wmt-site-content-strategy.scs_production.hp_summary_asset"
FYTD_START = "2026-01-31"

# ── ATF zone list ─────────────────────────────────────────────────────────────
ATF_ZONES_LIST = [
    'contentzone1','contentzone2','contentzone3',
    'contentzone4','contentzone5','contentzone6',
]
_ATF_SQL = "','".join(ATF_ZONES_LIST)   # ready for IN ('<zone>','<zone>',...)

# ── Module bucketing ──────────────────────────────────────────────────────────
# ATF  : HPOV | SIG | Item Carousel   (content_zone IN contentzone1–6)
# BTF  : Navigation | Content | Carousels | Utility  (all other zones)
#
# ⚠️  MM SIG LAUNCH — May 14, 2026 (Thursday night)
# Multi-Modal SIG changed ATF zone assignments:
#   BEFORE May 14 → HPOV in contentzone3, SIG in contentzone4/5/6
#   AFTER  May 14 → HPOV in contentzone2, SIG Top Picks in contentzone3,
#                   Rye + Site Merch carousels in contentzone4/5/6
# Bucketing below uses hp_module_name as primary key (unchanged by MM SIG),
# so classification remains correct across both periods. The ATF boundary
# (contentzone1–6) is also unchanged.
# Any metric spikes after May 15, 2026 → attribute to MM SIG first.
NEW_MODULE_CASE = f"""
  CASE
    -- HPOV: AutoScroll Cards 1-5
    -- Pre-MM SIG:  contentzone3
    -- Post-MM SIG (May 14, 2026+): contentzone2
    WHEN hp_module_name IN (
      'AutoScroll Card 1','AutoScroll Card 2','AutoScroll Card 3',
      'AutoScroll Card 4','AutoScroll Card 5'
    ) AND LOWER(content_zone) IN ('{_ATF_SQL}')
    THEN 'HPOV'

    -- SIG: SIG Cards 1-6
    -- Pre-MM SIG:  contentzone4/5/6
    -- Post-MM SIG: Top Picks → contentzone3; Rye/Site Merch → contentzone4/5/6
    WHEN hp_module_name IN (
      'SIG Card 1','SIG Card 2','SIG Card 3',
      'SIG Card 4','SIG Card 5','SIG Card 6'
    ) AND LOWER(content_zone) IN ('{_ATF_SQL}')
    THEN 'SIG'

    -- Item Carousel: moduletype = ItemCarousel, ATF zones only
    WHEN LOWER(COALESCE(moduletype,'')) = 'itemcarousel'
      AND LOWER(content_zone) IN ('{_ATF_SQL}')
    THEN 'Item Carousel'

    WHEN hp_module_type = 'Navigation'
      AND LOWER(content_zone) NOT IN ('{_ATF_SQL}')
    THEN 'Navigation'

    WHEN hp_module_type = 'Content'
      AND LOWER(content_zone) NOT IN ('{_ATF_SQL}')
    THEN 'Content'

    WHEN hp_module_type = 'Carousel'
      AND LOWER(content_zone) NOT IN ('{_ATF_SQL}')
    THEN 'Carousels'

    WHEN hp_module_type = 'Utility'
    THEN 'Utility'
  END
"""

# ── Canonical Merch content-type filter (content_served_by approach) ──────────
CONTENT_TYPE_FILTER = """
AND CASE
  WHEN LOWER(IFNULL(content_served_by,'')) = 'ads'
    THEN 'WMC'
  WHEN LOWER(IFNULL(disable_content_personalization,'')) LIKE '%true%'
    THEN 'Merch'
  WHEN LOWER(IFNULL(disable_content_personalization,'')) LIKE '%false%'
    AND LOWER(IFNULL(personalized_asset,'')) = 'default'
    AND session_start_dt <= '2025-03-01'
    AND LOWER(IFNULL(content_zone,'')) = 'contentzone3'
    AND LOWER(IFNULL(hp_module_name,'')) IN (
      'autoscroll card 1','autoscroll card 2','autoscroll card 3'
    )
    THEN 'WMC'
  WHEN LOWER(IFNULL(disable_content_personalization,'')) LIKE '%false%'
    AND LOWER(IFNULL(personalized_asset,'')) = 'default'
    AND (
      (
        LOWER(IFNULL(content_zone,'')) IN ('contentzone8','contentzone9')
        AND LOWER(IFNULL(hp_module_name,'')) = 'adjustable banner small'
      )
      OR
      (
        LOWER(IFNULL(content_zone,'')) IN ('contentzone10','contentzone11')
        AND LOWER(IFNULL(hp_module_name,'')) = 'triple pack small'
      )
    )
    THEN 'WMC'
  ELSE 'Merch'
END = 'Merch'
"""

MODULE_ORDER = ['HPOV','SIG','Item Carousel','Navigation','Content','Carousels','Utility']
ATF_MODULES  = {'HPOV','SIG','Item Carousel'}
BTF_MODULES  = {'Navigation','Content','Carousels','Utility'}


# ─────────────────────────────────────────────────────────────────────────────
#  Helpers
# ─────────────────────────────────────────────────────────────────────────────

def get_walmart_fiscal_week_dates(selected_date: str):
    """Saturday → selected_date (current) and same days one week prior."""
    dt      = datetime.strptime(selected_date, "%Y-%m-%d")
    weekday = dt.weekday()                         # Mon=0 … Sun=6
    days_since_saturday = 0 if weekday == 5 else (1 if weekday == 6 else weekday + 2)
    sat          = dt - timedelta(days=days_since_saturday)
    curr_start   = sat.strftime("%Y-%m-%d")
    curr_end     = dt.strftime("%Y-%m-%d")
    prev_start   = (sat - timedelta(days=7)).strftime("%Y-%m-%d")
    prev_end     = (dt  - timedelta(days=7)).strftime("%Y-%m-%d")
    return curr_start, curr_end, prev_start, prev_end


def fmt_pct(value, decimals=2):
    if value is None:
        return "—"
    return f"{round(value, decimals)}%"


def fmt_rate(value, decimals=1):
    """Format ATC Rate (no % sign, just the number)."""
    if value is None:
        return "—"
    return f"{round(value, decimals)}"


def fmt_wow(value):
    """Return (display_text, css_class) for a WoW value — no decimals."""
    if value is None:
        return ("—", "")
    rounded = int(round(value, 0))
    sign  = "+" if rounded > 0 else ""
    text  = f"{sign}{rounded}%"
    css   = "green" if rounded > 0 else ("red" if rounded < 0 else "")
    return (text, css)


# ─────────────────────────────────────────────────────────────────────────────
#  BigQuery: single query → FYTD baseline + current week
# ─────────────────────────────────────────────────────────────────────────────

def engagement_color(pct):
    """Return (hex_color, label) for engagement performance circle."""
    if pct is None:
        return ("#9ca3af", "—")
    if pct >= 100:
        return ("#16a34a", f"{pct}%")   # green
    if pct >= 90:
        return ("#d97706", f"{pct}%")   # amber
    return ("#dc2626", f"{pct}%")       # red


def get_module_insights(start_date: str, end_date: str,
                        prev_start: str, prev_end: str,
                        module_data: list = None) -> dict:
    """
    Run targeted deep-dives per module and return a dict of 1-liner insights.
    module_data is used to align carousel/message-level findings with the
    module-level WoW direction — so we never contradict what the summary shows.
    """
    # Build a lookup: module_name → row (gives us ctr_wow_pct, engagement_pct)
    mod_lookup = {r["module"]: r for r in (module_data or [])}
    insights   = {}

    # Helper: is module-level CTR trending up?
    def mod_up(name):
        wow = mod_lookup.get(name, {}).get("ctr_wow_pct")
        return wow is not None and wow > 0

    # ── HPOV: top CTR message + biggest WoW mover ────────────────────
    try:
        q = f"""
        WITH curr AS (
          SELECT message_name,
            ROUND(SAFE_DIVIDE(SUM(overall_click_count),SUM(module_view_count))*100,2) AS ctr
          FROM `{BQ_TABLE}`
          WHERE session_start_dt BETWEEN '{start_date}' AND '{end_date}'
            AND hp_module_name IN
              ('AutoScroll Card 1','AutoScroll Card 2','AutoScroll Card 3',
               'AutoScroll Card 4','AutoScroll Card 5')
            AND LOWER(content_zone) IN ('{_ATF_SQL}')
            {CONTENT_TYPE_FILTER}
          GROUP BY 1 HAVING SUM(module_view_count)>3000000
        ),
        prev AS (
          SELECT message_name,
            ROUND(SAFE_DIVIDE(SUM(overall_click_count),SUM(module_view_count))*100,2) AS ctr
          FROM `{BQ_TABLE}`
          WHERE session_start_dt BETWEEN '{prev_start}' AND '{prev_end}'
            AND hp_module_name IN
              ('AutoScroll Card 1','AutoScroll Card 2','AutoScroll Card 3',
               'AutoScroll Card 4','AutoScroll Card 5')
            AND LOWER(content_zone) IN ('{_ATF_SQL}')
            {CONTENT_TYPE_FILTER}
          GROUP BY 1 HAVING SUM(module_view_count)>3000000
        )
        SELECT c.message_name AS name,
          c.ctr AS ctr,
          p.ctr AS prev_ctr,
          ROUND((SAFE_DIVIDE(c.ctr,p.ctr)-1)*100,1) AS wow
        FROM curr c LEFT JOIN prev p USING(message_name)
        ORDER BY c.ctr DESC LIMIT 5
        """
        rows = list(client.query(q, job_config=JOB_CONFIG).result(timeout=45))
        if rows:
            valid = [r for r in rows if r.wow is not None]
            if mod_up("HPOV"):
                top = max(rows, key=lambda r: r.ctr)
                txt = f"{top.name} driving HPOV performance this week"
                gainers = [r for r in valid if r.wow > 0 and r.name != top.name]
                if gainers:
                    g2 = max(gainers, key=lambda r: r.wow)
                    txt += f"; {g2.name} also gaining WoW"
            else:
                losers = [r for r in valid if r.wow < 0]
                if losers:
                    l = min(losers, key=lambda r: r.wow)
                    txt = f"{l.name} pulling HPOV down this week"
                    if len(losers) > 1:
                        l2 = sorted(losers, key=lambda r: r.wow)[1]
                        txt += f"; {l2.name} also softening"
                else:
                    top = max(rows, key=lambda r: r.ctr)
                    txt = f"{top.name} holding steady WoW"
            insights["HPOV"] = txt
    except Exception as e:
        insights["HPOV"] = f"Query error: {str(e)[:60]}"

    # ── SIG: top carousel CTR + worst WoW mover ─────────────────────
    try:
        q = f"""
        WITH curr AS (
          SELECT COALESCE(NULLIF(carousel_name,''), hp_module_name) AS name,
            ROUND(SAFE_DIVIDE(SUM(overall_click_count),SUM(module_view_count))*100,2) AS ctr
          FROM `{BQ_TABLE}`
          WHERE session_start_dt BETWEEN '{start_date}' AND '{end_date}'
            AND hp_module_name IN
              ('SIG Card 1','SIG Card 2','SIG Card 3',
               'SIG Card 4','SIG Card 5','SIG Card 6')
            AND LOWER(content_zone) IN ('{_ATF_SQL}')
            AND lang_cd = 'English'
            {CONTENT_TYPE_FILTER}
          GROUP BY 1 HAVING SUM(module_view_count)>3000000
        ),
        prev AS (
          SELECT COALESCE(NULLIF(carousel_name,''), hp_module_name) AS name,
            ROUND(SAFE_DIVIDE(SUM(overall_click_count),SUM(module_view_count))*100,2) AS ctr
          FROM `{BQ_TABLE}`
          WHERE session_start_dt BETWEEN '{prev_start}' AND '{prev_end}'
            AND hp_module_name IN
              ('SIG Card 1','SIG Card 2','SIG Card 3',
               'SIG Card 4','SIG Card 5','SIG Card 6')
            AND LOWER(content_zone) IN ('{_ATF_SQL}')
            AND lang_cd = 'English'
            {CONTENT_TYPE_FILTER}
          GROUP BY 1 HAVING SUM(module_view_count)>3000000
        )
        SELECT c.name, c.ctr AS ctr, p.ctr AS prev_ctr,
          ROUND((SAFE_DIVIDE(c.ctr,p.ctr)-1)*100,1) AS wow
        FROM curr c LEFT JOIN prev p USING(name)
        ORDER BY c.ctr DESC LIMIT 5
        """
        rows = list(client.query(q, job_config=JOB_CONFIG).result(timeout=45))
        if rows:
            valid = [r for r in rows if r.wow is not None]
            if mod_up("SIG"):
                top = max(rows, key=lambda r: r.ctr)
                txt = f"{top.name} leading SIG this week"
                gainers = [r for r in valid if r.wow > 0 and r.name != top.name]
                if gainers:
                    g2 = max(gainers, key=lambda r: r.wow)
                    txt += f"; {g2.name} also gaining WoW"
            else:
                losers = [r for r in valid if r.wow < 0]
                if losers:
                    l = min(losers, key=lambda r: r.wow)
                    txt = f"{l.name} pulling SIG down this week"
                    if len(losers) > 1:
                        l2 = sorted(losers, key=lambda r: r.wow)[1]
                        txt += f"; {l2.name} also softening"
                else:
                    top = max(rows, key=lambda r: r.ctr)
                    txt = f"{top.name} holding steady WoW"
            insights["SIG"] = txt
    except Exception as e:
        insights["SIG"] = f"Query error: {str(e)[:60]}"

    # ── Item Carousel: overall CTR WoW driver ────────────────────────
    try:
        q = f"""
        WITH curr AS (
          SELECT COALESCE(NULLIF(carousel_name,''), hp_module_name) AS name,
            ROUND(SAFE_DIVIDE(SUM(overall_click_count),SUM(module_view_count))*100,2) AS ctr,
            SUM(total_atc_count) AS atc, SUM(module_view_count) AS imp
          FROM `{BQ_TABLE}`
          WHERE session_start_dt BETWEEN '{start_date}' AND '{end_date}'
            AND LOWER(COALESCE(moduletype,''))='itemcarousel'
            AND LOWER(content_zone) IN ('{_ATF_SQL}')
            {CONTENT_TYPE_FILTER}
          GROUP BY 1 HAVING imp>3000000
        ),
        prev AS (
          SELECT COALESCE(NULLIF(carousel_name,''), hp_module_name) AS name,
            ROUND(SAFE_DIVIDE(SUM(overall_click_count),SUM(module_view_count))*100,2) AS ctr
          FROM `{BQ_TABLE}`
          WHERE session_start_dt BETWEEN '{prev_start}' AND '{prev_end}'
            AND LOWER(COALESCE(moduletype,''))='itemcarousel'
            AND LOWER(content_zone) IN ('{_ATF_SQL}')
            {CONTENT_TYPE_FILTER}
          GROUP BY 1
        )
        SELECT c.name, c.ctr, p.ctr AS prev_ctr,
          ROUND((SAFE_DIVIDE(c.ctr,p.ctr)-1)*100,1) AS wow
        FROM curr c LEFT JOIN prev p USING(name)
        ORDER BY c.imp DESC LIMIT 3
        """
        rows = list(client.query(q, job_config=JOB_CONFIG).result(timeout=45))
        if rows:
            valid = [r for r in rows if r.wow is not None]
            module_is_up = mod_up("Item Carousel")

            if module_is_up:
                top = max(rows, key=lambda r: r.ctr)
                txt = f"{top.name} driving Item Carousel WoW lift"
                gainers = [r for r in valid if r.wow > 0 and r.name != top.name]
                if gainers:
                    g2 = max(gainers, key=lambda r: r.wow)
                    txt += f"; {g2.name} also gaining WoW"
            else:
                losers = [r for r in valid if r.wow < 0]
                if losers:
                    l = min(losers, key=lambda r: r.wow)
                    txt = f"{l.name} pulling Item Carousel down this week"
                    if len(losers) > 1:
                        l2 = sorted(losers, key=lambda r: r.wow)[1]
                        txt += f"; {l2.name} also softening"
                else:
                    top = max(rows, key=lambda r: r.ctr)
                    txt = f"{top.name} holding steady WoW"
            insights["Item Carousel"] = txt
    except Exception as e:
        insights["Item Carousel"] = f"Query error: {str(e)[:60]}"

    # ── BTF modules: top hp_module_name by CTR WoW ──────────────────
    BTF_TYPES = {
        "Navigation": "Navigation",
        "Content":    "Content",
        "Carousels":  "Carousel",
        "Utility":    "Utility",
    }
    for mod_key, hp_type in BTF_TYPES.items():
        try:
            zone_clause = (
                f"AND LOWER(content_zone) NOT IN ('{_ATF_SQL}')"
                if mod_key != "Utility" else ""
            )
            q = f"""
            WITH curr AS (
              SELECT
                CASE
                  WHEN '{hp_type}' = 'Content'
                  THEN COALESCE(NULLIF(message_name,''), hp_module_name)
                  ELSE hp_module_name
                END AS name,
                ROUND(SAFE_DIVIDE(SUM(overall_click_count),SUM(module_view_count))*100,2) AS ctr,
                SUM(module_view_count) AS imp
              FROM `{BQ_TABLE}`
              WHERE session_start_dt BETWEEN '{start_date}' AND '{end_date}'
                AND hp_module_type = '{hp_type}'
                {zone_clause}
                {CONTENT_TYPE_FILTER}
              GROUP BY 1 HAVING imp>3000000
            ),
            prev AS (
              SELECT
                CASE
                  WHEN '{hp_type}' = 'Content'
                  THEN COALESCE(NULLIF(message_name,''), hp_module_name)
                  ELSE hp_module_name
                END AS name,
                ROUND(SAFE_DIVIDE(SUM(overall_click_count),SUM(module_view_count))*100,2) AS ctr
              FROM `{BQ_TABLE}`
              WHERE session_start_dt BETWEEN '{prev_start}' AND '{prev_end}'
                AND hp_module_type = '{hp_type}'
                {zone_clause}
                {CONTENT_TYPE_FILTER}
              GROUP BY 1
            )
            SELECT c.name, c.ctr, p.ctr AS prev_ctr,
              ROUND((SAFE_DIVIDE(c.ctr,p.ctr)-1)*100,1) AS wow, c.imp
            FROM curr c LEFT JOIN prev p USING(name)
            ORDER BY ABS(COALESCE(ROUND((SAFE_DIVIDE(c.ctr,p.ctr)-1)*100,1),0)) DESC
            LIMIT 3
            """
            rows = list(client.query(q, job_config=JOB_CONFIG).result(timeout=45))
            if rows:
                valid = [r for r in rows if r.wow is not None]
                is_up = mod_up(mod_key)

                if is_up:
                    top = max(rows, key=lambda r: r.ctr)
                    txt = f"{top.name} leading {mod_key} this week"
                    gainers = [r for r in valid if r.wow > 0 and r.name != top.name]
                    if gainers:
                        g2 = max(gainers, key=lambda r: r.wow)
                        txt += f"; {g2.name} also gaining WoW"
                else:
                    losers = [r for r in valid if r.wow < 0]
                    if losers:
                        l = min(losers, key=lambda r: r.wow)
                        txt = f"{l.name} pulling {mod_key} down this week"
                        if len(losers) > 1:
                            l2 = sorted(losers, key=lambda r: r.wow)[1]
                            txt += f"; {l2.name} also softening"
                    else:
                        top = max(rows, key=lambda r: r.ctr)
                        txt = f"{top.name} holding steady WoW"
                insights[mod_key] = txt
            else:
                insights[mod_key] = "Insufficient data for this period"
        except Exception as e:
            insights[mod_key] = f"Query error: {str(e)[:60]}"

    return insights


def get_engagement_data(start_date: str, end_date: str,
                        prev_start: str, prev_end: str):
    """
    Returns one row per module with:
      ctr_baseline   – FYTD CTR
      ctr_pct        – current-week CTR
      ctr_wow_pct    – WoW vs PREVIOUS fiscal week
      atc_rate_baseline – FYTD ATC Rate (ATCs per 1000 impressions)
      atc_rate_pct      – current-week ATC Rate
      atc_rate_wow_pct  – WoW vs PREVIOUS fiscal week ATC Rate
      engagement_pct – CTR% / CTR Baseline × 100
    """
    query = f"""
    WITH
    -- ── FYTD baseline (Jan 31 2026 → latest date) ───────────────────────
    fytd AS (
      SELECT
        {NEW_MODULE_CASE} AS Module,
        SAFE_DIVIDE(SUM(overall_click_count), SUM(module_view_count)) AS fytd_ctr,
        SAFE_DIVIDE(SUM(total_atc_count), SUM(module_view_count)) * 1000 AS fytd_atc_rate,
        SUM(module_view_count) AS fytd_imp
      FROM `{BQ_TABLE}`
      WHERE session_start_dt BETWEEN '{FYTD_START}' AND (
          SELECT MAX(session_start_dt) FROM `{BQ_TABLE}`
          WHERE session_start_dt >= '{FYTD_START}'
        )
        {CONTENT_TYPE_FILTER}
      GROUP BY Module HAVING Module IS NOT NULL
    ),

    -- ── Current fiscal partial week ───────────────────────────────────────
    curr AS (
      SELECT
        {NEW_MODULE_CASE} AS Module,
        SAFE_DIVIDE(SUM(overall_click_count), SUM(module_view_count)) AS curr_ctr,
        SAFE_DIVIDE(SUM(total_atc_count), SUM(module_view_count)) * 1000 AS curr_atc_rate,
        SUM(module_view_count) AS curr_imp
      FROM `{BQ_TABLE}`
      WHERE session_start_dt BETWEEN '{start_date}' AND '{end_date}'
        {CONTENT_TYPE_FILTER}
      GROUP BY Module HAVING Module IS NOT NULL
    ),

    -- ── Previous fiscal week (WoW denominator) ────────────────────────────
    prev AS (
      SELECT
        {NEW_MODULE_CASE} AS Module,
        SAFE_DIVIDE(SUM(overall_click_count), SUM(module_view_count)) AS prev_ctr,
        SAFE_DIVIDE(SUM(total_atc_count), SUM(module_view_count)) * 1000 AS prev_atc_rate,
        SUM(module_view_count) AS prev_imp
      FROM `{BQ_TABLE}`
      WHERE session_start_dt BETWEEN '{prev_start}' AND '{prev_end}'
        {CONTENT_TYPE_FILTER}
      GROUP BY Module HAVING Module IS NOT NULL
    )

    SELECT
      c.Module,

      -- CTR Baseline = FYTD CTR
      ROUND(f.fytd_ctr * 100, 2) AS ctr_baseline,
      -- Current CTR %
      ROUND(c.curr_ctr * 100, 2) AS ctr_pct,
      -- CTR WoW = vs previous fiscal week
      ROUND((SAFE_DIVIDE(c.curr_ctr, p.prev_ctr) - 1) * 100, 1) AS ctr_wow_pct,

      -- ATC Rate Baseline = FYTD ATC Rate (ATCs per 1000 impressions)
      ROUND(f.fytd_atc_rate, 1) AS atc_rate_baseline,
      -- Current ATC Rate
      ROUND(c.curr_atc_rate, 1) AS atc_rate_pct,
      -- ATC Rate WoW = vs previous fiscal week
      ROUND((SAFE_DIVIDE(c.curr_atc_rate, p.prev_atc_rate) - 1) * 100, 1) AS atc_rate_wow_pct,

      -- Engagement Performance:
      -- CTR% / CTR Baseline × 100
      ROUND(SAFE_DIVIDE(c.curr_ctr, f.fytd_ctr) * 100, 0) AS engagement_pct

    FROM curr c
    LEFT JOIN fytd f  ON c.Module = f.Module
    LEFT JOIN prev p  ON c.Module = p.Module

    ORDER BY
      CASE c.Module
        WHEN 'HPOV'          THEN 1
        WHEN 'SIG'           THEN 2
        WHEN 'Item Carousel' THEN 3
        WHEN 'Navigation'    THEN 4
        WHEN 'Content'       THEN 5
        WHEN 'Carousels'     THEN 6
        WHEN 'Utility'       THEN 7
        ELSE 99
      END
    """
    try:
        rows = list(client.query(query, job_config=JOB_CONFIG).result(timeout=90))
        result = []
        for r in rows:
            row = {
                "module":         r.Module,
                "ctr_baseline":   r.ctr_baseline,
                "ctr_pct":        r.ctr_pct,
                "ctr_wow_pct":    r.ctr_wow_pct,
                "atc_rate_baseline":   r.atc_rate_baseline,
                "atc_rate_pct":        r.atc_rate_pct,
                "atc_rate_wow_pct":    r.atc_rate_wow_pct,
                "engagement_pct": r.engagement_pct,
            }
            # Utility: blank ATC Rate columns, but colour circle using CTR-only ratio
            if r.Module == "Utility":
                row["atc_rate_baseline"]   = None
                row["atc_rate_pct"]        = None
                row["atc_rate_wow_pct"]    = None
                if r.ctr_baseline and r.ctr_pct:
                    row["engagement_pct"] = round((r.ctr_pct / r.ctr_baseline) * 100, 0)
                else:
                    row["engagement_pct"] = None
            result.append(row)
        return result
    except Exception as e:
        print(f"[BQ Error] {e}")
        return []


# ─────────────────────────────────────────────────────────────────────────────
#  HTML table renderer  (merged ATF / BTF placement cells)
# ─────────────────────────────────────────────────────────────────────────────

def combined_signal(ctr_wow_pct, atc_rate_wow_pct, engagement_pct, module=None):
    """
    Reconcile WoW (vs prior week) with Engagement Performance (vs FYTD)
    into a single honest directional statement.
    """
    wow_up   = ctr_wow_pct is not None and ctr_wow_pct > 0
    wow_down = ctr_wow_pct is not None and ctr_wow_pct < 0
    ep_green = engagement_pct is not None and engagement_pct >= 100
    ep_amber = engagement_pct is not None and 90 <= engagement_pct < 100
    ep_red   = engagement_pct is not None and engagement_pct < 90

    # Both signals aligned — clearest cases
    if wow_up and ep_green:
        return "above FYTD & strengthening WoW — on track"
    if wow_down and ep_red:
        return "below FYTD & still falling WoW — needs attention"
    if wow_down and ep_amber:
        return "approaching FYTD lower bound — watch closely"

    # Divergent signals — explain the contradiction
    if wow_up and ep_red:
        return "bouncing WoW but well below FYTD baseline — recovery in early stages"
    if wow_up and ep_amber:
        return "recovering WoW, closing gap to FYTD — maintain momentum"
    if wow_down and ep_green:
        return "above FYTD baseline but easing WoW — monitor for continuation"

    return "tracking near FYTD baseline"


def render_engagement_table(data: list, insights: dict) -> str:
    if not data:
        return '<p style="color:#6b7280;font-style:italic;">No data returned from BigQuery. Check your credentials and date range.</p>'

    by_module = {r["module"]: r for r in data}
    ordered   = [by_module[m] for m in MODULE_ORDER if m in by_module]

    atf_rows = [r for r in ordered if r["module"] in ATF_MODULES]
    btf_rows = [r for r in ordered if r["module"] in BTF_MODULES]
    placement_done = {"ATF": False, "BTF": False}

    TH  = 'style="padding:10px 14px;text-align:center;font-weight:700;border:1px solid #1a3a60;white-space:nowrap;"'
    TH_L= 'style="padding:10px 14px;text-align:left;font-weight:700;border:1px solid #1a3a60;white-space:nowrap;"'
    TD  = 'style="padding:9px 14px;text-align:center;border:1px solid #d1d5db;"'
    TD_L= 'style="padding:9px 14px;text-align:left;border:1px solid #d1d5db;"'

    html = f'''
<table style="border-collapse:collapse;width:100%;font-family:\'Segoe UI\',Arial,sans-serif;font-size:13px;margin-top:4px;">
  <thead>
    <tr style="background:#041e42;color:white;">
      <th {TH}>Homepage Placement</th>
      <th {TH}>Module</th>
      <th {TH}>CTR Baseline</th>
      <th {TH}>CTR%</th>
      <th {TH}>WoW</th>
      <th {TH}>ATC Rate Baseline</th>
      <th {TH}>ATC Rate</th>
      <th {TH}>WoW</th>
      <th {TH}>Engagement Performance</th>
      <th {TH}>% to Engagement</th>
      <th {TH_L}>Insights</th>
    </tr>
  </thead>
  <tbody>
'''
    COLOR = {"green": "#16a34a", "red": "#dc2626", "": "#1a1a1a"}

    for row in ordered:
        mod       = row["module"]
        placement = "ATF" if mod in ATF_MODULES else "BTF"
        group     = atf_rows if placement == "ATF" else btf_rows

        ctr_wow_txt, ctr_wow_css = fmt_wow(row["ctr_wow_pct"])
        atc_wow_txt, atc_wow_css = fmt_wow(row["atc_rate_wow_pct"])
        ctr_color = COLOR[ctr_wow_css]
        atc_color = COLOR[atc_wow_css]

        # Build full insight = BQ detail + reconciled WoW-vs-EP signal
        bq_detail  = insights.get(mod, "")
        signal     = combined_signal(
            row["ctr_wow_pct"], row["atc_rate_wow_pct"],
            row["engagement_pct"], module=mod
        )
        full_insight = f"{bq_detail}. {signal}." if bq_detail else f"{signal}."

        # Engagement Performance — circle only (no number)
        eng_color, _ = engagement_color(row["engagement_pct"])
        circle_only = (
            f'<div style="display:flex;align-items:center;justify-content:center;">'
            f'<span style="display:inline-block;width:24px;height:24px;border-radius:50%;'
            f'background:{eng_color};"></span></div>'
        )

        tr = '<tr>'

        # ── Homepage Placement cell (merged) ──
        if not placement_done[placement]:
            placement_done[placement] = True
            rowspan = len(group)
            tr += (
                f'<td rowspan="{rowspan}" '
                f'style="padding:10px 14px;text-align:center;font-weight:800;'
                f'font-size:16px;border:1px solid #d1d5db;background:#ffffff;'
                f'vertical-align:middle;color:#041e42;">{placement}</td>'
            )

        tr += f'''
      <td {TD}><strong>{mod}</strong></td>
      <td {TD}>{fmt_pct(row["ctr_baseline"])}</td>
      <td {TD}>{fmt_pct(row["ctr_pct"])}</td>
      <td style="padding:9px 14px;text-align:center;border:1px solid #d1d5db;font-weight:700;color:{ctr_color};">{ctr_wow_txt}</td>
      <td {TD}>{fmt_rate(row["atc_rate_baseline"])}</td>
      <td {TD}>{fmt_rate(row["atc_rate_pct"])}</td>
      <td style="padding:9px 14px;text-align:center;border:1px solid #d1d5db;font-weight:700;color:{atc_color};">{atc_wow_txt}</td>
      <td {TD}>{circle_only}</td>
      <td {TD}>{fmt_pct(row["engagement_pct"], 0)}</td>
      <td {TD_L} style="padding:9px 14px;text-align:left;border:1px solid #d1d5db;font-size:12px;color:#374151;line-height:1.5;">{full_insight}</td>
    </tr>'''
        html += tr

    html += "  </tbody>\n</table>"
    return html


# ─────────────────────────────────────────────────────────────────────────────
#  PPT generator
# ─────────────────────────────────────────────────────────────────────────────

def generate_ppt_bytes(data: list, current_start: str, current_end: str, insights: dict) -> BytesIO:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

    NAVY  = RGBColor(4,   30,  66)
    WHITE = RGBColor(255, 255, 255)
    GREEN = RGBColor(22,  163, 74)
    RED   = RGBColor(220, 38,  38)
    BLACK = RGBColor(0,   0,   0)
    FONT  = "Everyday Sans"
    FSIZE = Pt(11)
    BORDER_COLOR = "000000"
    BORDER_W     = 12700   # 1pt in EMUs

    def set_borders(cell, color=BORDER_COLOR, w=BORDER_W, anchor='ctr'):
        """Add solid borders + set vertical alignment on a table cell via XML.
        anchor: 't'=top  'ctr'=middle  'b'=bottom
        """
        tc   = cell._tc
        tcPr = tc.find(qn('a:tcPr'))
        if tcPr is None:
            tcPr = etree.SubElement(tc, qn('a:tcPr'))
        # vertical anchor — must be set on tcPr, not text_frame
        tcPr.set('anchor', anchor)
        for side in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
            tag = qn(side)
            for old in tcPr.findall(tag):
                tcPr.remove(old)
            ln = etree.SubElement(tcPr, tag,
                                  w=str(w), cap='flat', cmpd='sng', algn='ctr')
            sf = etree.SubElement(ln, qn('a:solidFill'))
            etree.SubElement(sf, qn('a:srgbClr'), val=color)
            etree.SubElement(ln, qn('a:prstDash'), val='solid')

    def style_cell(cell, text='', bold=False, font_color=BLACK,
                   bg=WHITE, align=PP_ALIGN.CENTER, font_size=FSIZE,
                   wrap=True):
        """Set fill, text, font, alignment and borders on a cell."""
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.text      = text
        p.alignment = align
        p.font.name  = FONT
        p.font.size  = font_size
        p.font.bold  = bold
        p.font.color.rgb = font_color
        set_borders(cell)  # sets anchor='ctr' (middle) + all four borders

    # ── Presentation setup ─────────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width  = Inches(16)
    prs.slide_height = Inches(9)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    tb = slide.shapes.add_textbox(Inches(0.35), Inches(0.10), Inches(15.3), Inches(0.45))
    p  = tb.text_frame.paragraphs[0]
    p.text = "Homepage Engagement Performance"
    p.font.name = FONT; p.font.size = Pt(18); p.font.bold = True
    p.font.color.rgb = NAVY

    # Subtitle
    tb2 = slide.shapes.add_textbox(Inches(0.35), Inches(0.55), Inches(15.3), Inches(0.3))
    p2  = tb2.text_frame.paragraphs[0]
    p2.text = (f"Week: {current_start} \u2192 {current_end}  |  "
               f"CTR Baseline = FYTD (Jan 31 2026 \u2192 latest)  |  "
               f"ATC Rate = ATCs per 1000 impressions  |  Merch only")
    p2.font.name = FONT; p2.font.size = Pt(9)
    p2.font.color.rgb = RGBColor(100, 100, 100)

    # ── Build ordered data ─────────────────────────────────────────────────
    by_module = {r["module"]: r for r in data}
    ordered   = [by_module[m] for m in MODULE_ORDER if m in by_module]
    n_rows    = len(ordered) + 1
    n_cols    = 11   # added % to Engagement column
    col_w     = [1.05, 1.35, 1.25, 1.1, 0.9, 1.25, 1.1, 0.9, 1.0, 1.1, 4.25]

    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.35), Inches(0.90),
        Inches(sum(col_w)), Inches(7.85)
    ).table
    for i, w in enumerate(col_w):
        tbl.columns[i].width = Inches(w)

    # ── Header row ─────────────────────────────────────────────────────────
    headers = [
        "Homepage\nPlacement", "Module",
        "CTR Baseline", "CTR%", "WoW",
        "ATC Rate\nBaseline", "ATC Rate", "WoW",
        "Engagement\nPerformance", "% to\nEngagement", "Insights"
    ]
    for j, h in enumerate(headers):
        style_cell(tbl.cell(0, j), text=h, bold=True,
                   font_color=WHITE, bg=NAVY,
                   align=PP_ALIGN.CENTER, font_size=FSIZE)

    # ── Data rows — col 0 left blank for now (filled after merge) ──────────
    atf_indices = [i+1 for i, r in enumerate(ordered) if r["module"] in ATF_MODULES]
    btf_indices = [i+1 for i, r in enumerate(ordered) if r["module"] in BTF_MODULES]

    def wow_str(v):
        if v is None: return "—"
        rounded = int(round(v, 0))
        return f"{'+'if rounded>0 else ''}{rounded}%"

    for i, row in enumerate(ordered):
        ri  = i + 1
        mod = row["module"]

        # col 0 — leave empty; placement text set after merge
        style_cell(tbl.cell(ri, 0), text='')

        ppt_signal = combined_signal(
            row["ctr_wow_pct"], row["atc_rate_wow_pct"],
            row["engagement_pct"], module=mod
        )
        ppt_bq = insights.get(mod, "")
        ppt_insight = f"{ppt_bq}. {ppt_signal}." if ppt_bq else f"{ppt_signal}."

        data_vals = [
            mod,
            fmt_pct(row["ctr_baseline"]),
            fmt_pct(row["ctr_pct"]),
            wow_str(row["ctr_wow_pct"]),
            fmt_rate(row["atc_rate_baseline"]),
            fmt_rate(row["atc_rate_pct"]),
            wow_str(row["atc_rate_wow_pct"]),
            "\u25cf",
            fmt_pct(row["engagement_pct"], 0),
            ppt_insight,                        # col 10 — full insight
        ]
        for j, val in enumerate(data_vals):
            col = j + 1
            align = PP_ALIGN.LEFT if col == 10 else PP_ALIGN.CENTER
            style_cell(tbl.cell(ri, col), text=val, align=align)

            # CTR WoW colour (col 4)
            if col == 4 and row["ctr_wow_pct"] is not None:
                p = tbl.cell(ri, col).text_frame.paragraphs[0]
                p.font.bold = True
                p.font.color.rgb = GREEN if row["ctr_wow_pct"] > 0 else (RED if row["ctr_wow_pct"] < 0 else BLACK)
            # ATC Rate WoW colour (col 7)
            elif col == 7 and row["atc_rate_wow_pct"] is not None:
                p = tbl.cell(ri, col).text_frame.paragraphs[0]
                p.font.bold = True
                p.font.color.rgb = GREEN if row["atc_rate_wow_pct"] > 0 else (RED if row["atc_rate_wow_pct"] < 0 else BLACK)
            # Engagement Performance circle colour (col 8) — large ● char
            elif col == 8:
                ep = row["engagement_pct"]
                ep_rgb = (GREEN if ep is not None and ep >= 100
                          else (RGBColor(217, 119, 6) if ep is not None and ep >= 90
                          else RED))
                p = tbl.cell(ri, col).text_frame.paragraphs[0]
                p.font.size  = Pt(20)   # make the dot visually large
                p.font.bold  = False
                p.font.color.rgb = ep_rgb if ep is not None else RGBColor(156, 163, 175)
            # % to Engagement colour (col 9)
            elif col == 9:
                ep = row["engagement_pct"]
                if ep is not None:
                    ep_rgb = (GREEN if ep >= 100
                              else (RGBColor(217, 119, 6) if ep >= 90
                              else RED))
                    p = tbl.cell(ri, col).text_frame.paragraphs[0]
                    p.font.bold = True
                    p.font.color.rgb = ep_rgb

    # ── Merge col 0 for ATF and BTF, THEN write text once ─────────────────
    for idx_list, label in [(atf_indices, "ATF"), (btf_indices, "BTF")]:
        if not idx_list:
            continue
        if len(idx_list) >= 2:
            tbl.cell(idx_list[0], 0).merge(tbl.cell(idx_list[-1], 0))
        # Only now write text — after merge so no duplication
        mc = tbl.cell(idx_list[0], 0)
        style_cell(mc, text=label, bold=True, font_color=NAVY,
                   bg=WHITE, align=PP_ALIGN.CENTER, font_size=Pt(14))
        set_borders(mc)   # re-apply borders after merge rewrites tcPr

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ─────────────────────────────────────────────────────────────────────────────
#  Routes
# ─────────────────────────────────────────────────────────────────────────────

# ── Simple in-memory cache keyed by selected_date ────────────────────────────
# Avoids re-running all BQ queries when user clicks Download PPT immediately
# after loading the page.
_CACHE: dict = {}


@app.get("/", response_class=HTMLResponse)
async def home(selected_date: str = Query(default=datetime.today().strftime("%Y-%m-%d"))):
    curr_start, curr_end, prev_start, prev_end = get_walmart_fiscal_week_dates(selected_date)
    data       = get_engagement_data(curr_start, curr_end, prev_start, prev_end)
    insights   = get_module_insights(curr_start, curr_end, prev_start, prev_end, module_data=data)
    # Store in cache so PPT download can reuse without re-querying BQ
    _CACHE[selected_date] = (data, insights, curr_start, curr_end)
    table_html = render_engagement_table(data, insights)

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width,initial-scale=1.0">
  <title>WBR — Homepage Engagement Performance</title>
  <style>
    *   {{ box-sizing:border-box; margin:0; padding:0; }}
    body{{ background:#f0f2f5; font-family:'Segoe UI',Arial,sans-serif; padding:24px; }}
    .card{{ background:white; border-radius:12px; box-shadow:0 2px 14px rgba(0,0,0,.09); padding:30px; max-width:1500px; margin:0 auto; }}

    /* ── top bar ── */
    .topbar{{ display:flex; align-items:flex-start; justify-content:space-between; margin-bottom:22px; gap:16px; flex-wrap:wrap; }}
    .title  {{ font-size:22px; font-weight:800; color:#041e42; }}
    .subtitle{{ font-size:12px; color:#6b7280; margin-top:5px; line-height:1.6; }}

    /* ── date form ── */
    .datebar{{ background:#f8f9fa; border-radius:8px; padding:14px 18px; margin-bottom:22px;
               display:flex; align-items:center; gap:14px; flex-wrap:wrap; border:1px solid #e5e7eb; }}
    .datebar label{{ font-size:13px; font-weight:600; color:#374151; }}
    .datebar input {{ border:1px solid #d1d5db; border-radius:6px; padding:7px 11px; font-size:13px; }}
    .datebar button{{ background:#041e42; color:white; border:none; border-radius:6px;
                      padding:8px 20px; font-size:13px; cursor:pointer; font-weight:600; }}
    .datebar button:hover{{ background:#0a3161; }}

    /* ── week badges ── */
    .badges{{ display:flex; gap:10px; flex-wrap:wrap; }}
    .badge {{ padding:5px 13px; border-radius:6px; font-size:12px; font-weight:600; white-space:nowrap; }}
    .badge-c{{ background:#dce8f5; border:1px solid #041e42; color:#041e42; }}
    .badge-p{{ background:#f3f4f6; border:1px solid #9ca3af; color:#6b7280; }}

    /* ── download button ── */
    .dl-btn{{ background:#0071ce; color:white; border:none; border-radius:8px;
              padding:10px 22px; font-size:14px; cursor:pointer; font-weight:700;
              text-decoration:none; display:inline-flex; align-items:center; gap:8px;
              white-space:nowrap; }}
    .dl-btn:hover{{ background:#005ba0; }}

    /* ── legend / note ── */
    .legend{{ margin-top:12px; font-size:12px; color:#6b7280; }}
    .g{{ color:#16a34a; font-weight:700; }}
    .r{{ color:#dc2626; font-weight:700; }}
    .note{{ font-size:11px; color:#9ca3af; margin-top:12px; padding-top:12px; border-top:1px solid #f0f2f5; line-height:1.8; }}
    code{{ background:#f3f4f6; padding:1px 5px; border-radius:3px; font-size:11px; }}
  </style>
</head>
<body>
<div class="card">

  <!-- ── Top bar ─────────────────────────────────────────── -->
  <div class="topbar">
    <div>
      <div class="title">📊 Homepage Engagement Performance</div>
      <div class="subtitle">
        CTR Baseline = FYTD (Jan 31 2026 → latest available date) &nbsp;|&nbsp;
        ATC Rate = ATCs per 1000 impressions &nbsp;|&nbsp;
        WoW = Current partial week vs FYTD baseline<br>
        ATF = contentZone1–6 &nbsp;|&nbsp; BTF = all other zones &nbsp;|&nbsp; Merch only
      </div>
    </div>
    <a class="dl-btn" href="/download-ppt?selected_date={selected_date}">⬇&nbsp; Download PPT</a>
  </div>

  <!-- ── Date picker ─────────────────────────────────────── -->
  <div class="datebar">
    <form method="get" style="display:flex;align-items:center;gap:12px;flex-wrap:wrap;">
      <label>Select Date (Fiscal Week End)</label>
      <input type="date" name="selected_date" value="{selected_date}">
      <button type="submit">Generate Report</button>
    </form>
    <div class="badges">
      <span class="badge badge-c">📅 Current: {curr_start} → {curr_end}</span>
      <span class="badge badge-p">🔙 Prior: {prev_start} → {prev_end}</span>
    </div>
  </div>

  <!-- ── Table ───────────────────────────────────────────── -->
  <div style="overflow-x:auto;">
    {table_html}
  </div>

  <!-- ── Legend ──────────────────────────────────────────── -->
  <div class="legend">
    <span class="g">■ Green WoW</span> = current week above FYTD baseline &nbsp;&nbsp;
    <span class="r">■ Red WoW</span> = current week below FYTD baseline
  </div>

  <div class="note">
    Source: <code>hp_summary_asset</code> &nbsp;|&nbsp;
    Content Type: Merch only (canonical <code>content_served_by</code> CASE logic) &nbsp;|&nbsp;
    WoW formula: (Current / FYTD − 1) × 100
  </div>

</div>
</body>
</html>"""
    return HTMLResponse(content=html)


@app.get("/download-ppt")
async def download_ppt(selected_date: str = Query(default=datetime.today().strftime("%Y-%m-%d"))):
    try:
        # Reuse cached data if the user already loaded the page for this date
        if selected_date in _CACHE:
            data, insights, curr_start, curr_end = _CACHE[selected_date]
        else:
            curr_start, curr_end, prev_start, prev_end = get_walmart_fiscal_week_dates(selected_date)
            data     = get_engagement_data(curr_start, curr_end, prev_start, prev_end)
            insights = get_module_insights(curr_start, curr_end, prev_start, prev_end, module_data=data)
            _CACHE[selected_date] = (data, insights, curr_start, curr_end)

        ppt_buf  = generate_ppt_bytes(data, curr_start, curr_end, insights)
        filename = f"wbr-engagement-{curr_start}-to-{curr_end}.pptx"
        return StreamingResponse(
            ppt_buf,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
    except Exception as e:
        import traceback
        print(f"[PPT ERROR] {traceback.format_exc()}")
        from fastapi.responses import JSONResponse
        return JSONResponse(status_code=500, content={"error": str(e)})


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8001)
